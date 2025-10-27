import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from llm.chatgpt import ChatGpt
from random import choice
import google.generativeai as genai



def generate_powerpoint(topic: str):
    """
    Generate a visually enhanced PowerPoint presentation on a given topic.
    Uses Gemini if available, otherwise builds detailed fallback content.
    """

    try:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        output_dir = os.path.join(base_dir, "GeneratedPresentations")
        os.makedirs(output_dir, exist_ok=True)

        prompt = f"""
        Create a detailed PowerPoint outline for the topic "{topic}".
        Return valid JSON only.
        Format:
        {{
          "slides": [
            {{
              "title": "Introduction to {topic}",
              "content": ["...", "...", "..."]
            }},
            ...
          ]
        }}
        Include at least 7 slides with strong educational content, 
        interesting facts, and clear structure. 
        Avoid markdown or formatting, just plain JSON.
        """

        response = ChatGpt(prompt)
        slides_data = []

        if response and "slides" in response:
            try:
                slides_data = json.loads(response).get("slides", [])
            except Exception:
                import re
                json_match = re.search(r"\{[\s\S]*\}", response)
                if json_match:
                    try:
                        data = json.loads(json_match.group())
                        slides_data = data.get("slides", [])
                    except Exception:
                        pass

        if not slides_data:
            print("[PowerPointer] ⚠️ Gemini failed, using offline fallback.")
            slides_data = [
                {"title": f"Introduction to {topic}", "content": [
                    f"An overview of {topic}.",
                    f"Why {topic} is important in the modern world.",
                    "How it connects to everyday life."
                ]},
                {"title": f"Background of {topic}", "content": [
                    f"Origin and history of {topic}.",
                    "How it evolved over time.",
                    "Key contributors or discoveries."
                ]},
                {"title": "Core Concepts", "content": [
                    "Main principles or mechanisms.",
                    "Scientific or technical foundations.",
                    "Important terminology."
                ]},
                {"title": "Applications", "content": [
                    f"How {topic} is used in daily life.",
                    f"Key industries benefiting from {topic}.",
                    "Examples of real-world use cases."
                ]},
                {"title": "Challenges & Limitations", "content": [
                    f"Current issues with {topic}.",
                    "Research gaps and controversies.",
                    "Technical and ethical considerations."
                ]},
                {"title": "Interesting Facts", "content": [
                    f"Did you know? Fascinating facts about {topic}.",
                    f"Unexpected uses of {topic}.",
                    "Surprising statistics or trivia."
                ]},
                {"title": "Conclusion", "content": [
                    f"Summary of what we learned about {topic}.",
                    "Future outlook and possibilities.",
                    "Thank you!"
                ]}
            ]

        # Create Presentation
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]

        # Theme colors
        bg_colors = [
            RGBColor(32, 55, 100),   # navy
            RGBColor(45, 60, 85),    # blue-gray
            RGBColor(20, 40, 70),    # deep blue
            RGBColor(15, 30, 50),    # dark slate
        ]
        text_color = RGBColor(255, 255, 255)

        # Title Slide
        slide = prs.slides.add_slide(title_slide_layout)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = choice(bg_colors)
        slide.shapes.title.text = f"{topic.title()}"
        slide.placeholders[1].text = "A Presentation by Jarvis Assistant"
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = text_color
                        run.font.size = Pt(40 if shape == slide.shapes.title else 20)

        # Content Slides
        for slide_data in slides_data:
            slide = prs.slides.add_slide(content_layout)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = choice(bg_colors)

            title = slide.shapes.title
            title.text = slide_data["title"]
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = text_color

            content = slide.placeholders[1].text_frame
            content.clear()
            for bullet in slide_data["content"]:
                p = content.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = text_color

        # Save
        filename = f"{topic.replace(' ', '_')}.pptx"
        file_path = os.path.join(output_dir, filename)
        prs.save(file_path)

        print(f"[PowerPointer] ✅ PowerPoint saved: {file_path}")
        return file_path

    except Exception as e:
        print("[PowerPointer] ❌ Error:", e)
        return None
